import os
import shutil
import fitz
from fitz import EmptyFileError
from PIL import Image
import pytesseract
import io
from pptx import Presentation
from docx import Document
from docx.opc.exceptions import PackageNotFoundError
from odf.opendocument import load
from odf.text import P, Span
import openpyxl
import nbformat
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from tqdm import tqdm
from .config import EMBEDDING_MODEL, FAISS_INDEX_DIR, DB_NAME
from .database import store_chunks, init_db, get_total_chunks, fetch_chunks_batch

def extract_pdf_text(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        pdf_text = ''

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            if text.strip():
                pdf_text += text + '\n'
            else:
                print(f"[OCR] Page {page_num} contains no selectable text. Running OCR explicitly.")
                images = page.get_images(full=True)
                for img_index, img in enumerate(images, start=1):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image = Image.open(io.BytesIO(base_image["image"]))
                    ocr_text = pytesseract.image_to_string(image, lang='eng', config='--psm 6')
                    pdf_text += ocr_text + '\n'
                    print(f"[OCR] Extracted text explicitly from image {img_index} on page {page_num}")

        return pdf_text

    except EmptyFileError as e:
        print(f"⚠️ Skipping empty or corrupt file '{pdf_path}': {e}")
        return ''

def extract_python_text(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_notebook_text(notebook_path):
    notebook = nbformat.read(notebook_path, as_version=4)
    text = ''
    for cell in notebook.cells:
        if cell.cell_type == 'code':
            text += cell.source + '\n\n'
        elif cell.cell_type == 'markdown':
            text += cell.source + '\n\n'
    return text

def extract_text_from_pdf_images(pdf_path):
    doc = fitz.open(pdf_path)
    full_text = ''

    for page_num, page in enumerate(doc):
        images = page.get_images(full=True)
        print(f"Found {len(images)} images on page {page_num + 1}")

        for img_index, img in enumerate(images, start=1):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]

            image = Image.open(io.BytesIO(image_bytes))

            # OCR explicitly with PyTesseract
            ocr_text = pytesseract.image_to_string(image)
            full_text += ocr_text + '\n'

    return full_text

def extract_pptx_text(pptx_path):
    prs = Presentation(pptx_path)
    full_text = ''
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = ''
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + '\n'
        full_text += slide_text + '\n'
    return full_text


def extract_docx_text(docx_path):
    try:
        doc = Document(docx_path)
        return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    except PackageNotFoundError as e:
        print(f"⚠️ Skipping file '{docx_path}': {e}")
        return ''
    
def extract_odt_text(odt_path):
    doc = load(odt_path)
    paragraphs = doc.getElementsByType(P)

    def extract_text_from_element(element):
        text_content = ""
        for node in element.childNodes:
            if node.nodeType == node.TEXT_NODE:
                text_content += str(node.data)
            elif node.nodeType == node.ELEMENT_NODE and node.tagName == "text:span":
                text_content += extract_text_from_element(node)
        return text_content

    text = '\n'.join(extract_text_from_element(p) for p in paragraphs)
    return text

def extract_xlsx_text(xlsx_path):
    wb = openpyxl.load_workbook(xlsx_path, read_only=True)
    text = ''
    for sheet in wb.worksheets:
        text += f"Sheet: {sheet.title}\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
            text += row_text + '\n'
    return text

def chunk_text(text, chunk_size=500):
    words = text.split()
    return [' '.join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def is_excluded(path, exclusions):
    path_lower = path.lower()
    exclusions_lower = [ex.lower() for ex in exclusions]
    return any(exclusion in path_lower for exclusion in exclusions_lower)

def clear_previous_data(clear_faiss=True, clear_db=True):
    if clear_faiss and os.path.exists(FAISS_INDEX_DIR):
        shutil.rmtree(FAISS_INDEX_DIR)
        print(f"🗑️ Deleted previous FAISS index at {FAISS_INDEX_DIR}")

    if clear_db and os.path.exists(DB_NAME):
        os.remove(DB_NAME)
        print(f"🗑️ Deleted previous SQLite database at {DB_NAME}")

def process_files(root_folders, chunk_size=500, exclusions=None):
    if exclusions is None:
        exclusions = []

    chunks, sources = [], []

    with init_db() as conn:  # ensure proper closure explicitly
        for root_folder in root_folders:
            print(f"🔍 Starting recursive iteration from root folder: {root_folder}")

            for dirpath, dirnames, filenames in os.walk(root_folder):
                if is_excluded(dirpath, exclusions):
                    print(f"⛔️ Excluding directory: {dirpath}")
                    dirnames[:] = []
                    continue

                print(f"📂 Iterating directory: {dirpath}")

                for filename in filenames:
                    file_path = os.path.join(dirpath, filename)

                    if is_excluded(file_path, exclusions):
                        print(f"⛔️ Excluding file: {file_path}")
                        continue

                    file_ext = filename.lower()

                    try:
                        if file_ext.endswith('.pdf'):
                            print(f"📄 Processing PDF: {file_path}")
                            text = extract_pdf_text(file_path)
                        elif file_ext.endswith('.py'):
                            print(f"🐍 Processing Python file: {file_path}")
                            text = extract_python_text(file_path)
                        elif file_ext.endswith('.ipynb'):
                            print(f"📓 Processing Notebook: {file_path}")
                            text = extract_notebook_text(file_path)
                        elif file_ext.endswith('.pptx'):
                            print(f"📽️ Processing PowerPoint file: {file_path}")
                            text = extract_pptx_text(file_path)
                        elif file_ext.endswith('.docx'):
                            print(f"📃 Processing Word document: {file_path}")
                            text = extract_docx_text(file_path)
                        elif file_ext.endswith('.odt'):
                            print(f"📝 Processing ODT file: {file_path}")
                            text = extract_odt_text(file_path)
                        elif file_ext.endswith('.xlsx'):
                            print(f"📊 Processing Excel file: {file_path}")
                            text = extract_xlsx_text(file_path)
                        else:
                            continue
                    except Exception as e:
                        print(f"⚠️ Skipping file '{file_path}': {e}")
                        continue

                    file_chunks = chunk_text(text, chunk_size=chunk_size)
                    chunks.extend(file_chunks)
                    sources.extend([file_path] * len(file_chunks))
                    store_chunks(conn, file_chunks, [file_path] * len(file_chunks))

    print(f"✅ Finished processing files. Total chunks created: {len(chunks)}")
    return chunks, sources

def build_faiss_index(
    root_folders=None, 
    batch_size=100, 
    chunk_size=500, 
    exclusions=None, 
    rebuild_chunks=True
):
    embeddings_model = OpenAIEmbeddings(model=EMBEDDING_MODEL)

    if rebuild_chunks:
        if root_folders is None:
            raise ValueError("root_folders must be provided when rebuilding chunks.")
        
        clear_previous_data(clear_faiss=True, clear_db=True)
        process_files(root_folders, chunk_size=chunk_size, exclusions=exclusions)
    else:
        clear_previous_data(clear_faiss=True, clear_db=False)
        print("🔄 Skipping chunk extraction, using existing chunks in database.")

    with init_db() as conn:  # explicitly correct use of context manager
        total_chunks = get_total_chunks(conn)

        if total_chunks == 0:
            raise ValueError("No chunks found in the database. Cannot build embeddings.")

        vectorstore = None

        for offset in tqdm(range(0, total_chunks, batch_size), desc="🔧 Creating embeddings"):
            batch = fetch_chunks_batch(conn, offset, batch_size)
            batch_chunks, batch_sources = zip(*batch)
            batch_meta = [{"source": source} for source in batch_sources]

            if vectorstore is None:
                vectorstore = FAISS.from_texts(batch_chunks, embeddings_model, batch_meta)
            else:
                vectorstore.add_texts(batch_chunks, batch_meta)

        vectorstore.save_local(FAISS_INDEX_DIR)

    print(f"🎉 Embeddings built and saved to {FAISS_INDEX_DIR}")
